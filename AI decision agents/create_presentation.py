#!/usr/bin/env python3
"""
Generate PowerPoint presentation for AI Agent Platform Research
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (professional blue/gray)
    PRIMARY_COLOR = RGBColor(0, 51, 102)  # Dark blue
    ACCENT_COLOR = RGBColor(0, 112, 192)  # Light blue
    GRAY_COLOR = RGBColor(89, 89, 89)     # Dark gray

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = GRAY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_COLOR
            p.space_before = Pt(8)
            p.level = 0

        return slide

    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # Add table
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3),
                                       Inches(9), Inches(5.5)).table

        # Set column widths
        col_width = int(Inches(9) / num_cols)
        for col_idx in range(num_cols):
            table.columns[col_idx].width = col_width

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            text_frame = cell.text_frame
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                text_frame = cell.text_frame
                text_frame.paragraphs[0].font.size = Pt(11)
                text_frame.paragraphs[0].font.color.rgb = GRAY_COLOR

        return slide

    # Slide 1: Title
    add_title_slide(
        "AI Agent Platform Research",
        "Decision Framework for Government Benefits Adjudication"
    )

    # Slide 2: Executive Summary
    add_content_slide(
        "Executive Summary",
        [
            "• Goal: Evaluate AI agent platforms and tools to accelerate benefits adjudication decisions",
            "",
            "• Key Question: Buy tools, subscribe to platform/SaaS, or build custom?",
            "",
            "• Scope: 25+ platforms across enterprise, low-code, and code frameworks",
            "",
            "• Recommendation: Hybrid approach - enterprise platform + custom AI components",
            "",
            "• Timeline: 3-6 months for pilot, 9-12 months to production",
            "",
            "• Investment: $1M-2M over 3 years (platform-based approach)"
        ]
    )

    # Slide 3: Key Decision Factors
    add_content_slide(
        "Critical Requirements for Government Benefits",
        [
            "1. Compliance & Security",
            "   • FedRAMP High authorization required",
            "   • HIPAA compliance for health-related benefits",
            "   • SOC 2 Type II, ISO 27001 certifications",
            "",
            "2. Explainability & Auditability",
            "   • All decisions must have human-readable explanations",
            "   • Complete audit trails for appeals process",
            "",
            "3. Integration Capabilities",
            "   • Connect to legacy case management systems",
            "   • Access to multiple databases and document repositories",
            "",
            "4. Customization for Complex Rules",
            "   • Benefits eligibility varies by jurisdiction",
            "   • Multi-step adjudication workflows"
        ]
    )

    # Slide 4: Platform Landscape Overview
    add_content_slide(
        "Platform Categories",
        [
            "1. Enterprise AI Agent Platforms (Full-Stack SaaS)",
            "   • Pega, Appian, Azure AI Studio, Vertex AI, AWS Bedrock",
            "   • Complete solutions: workflow + case management + AI",
            "   • Fastest time to value (3-6 months)",
            "",
            "2. Low-Code/No-Code Builders",
            "   • Retool, Flowise, Dust, n8n, Voiceflow",
            "   • Rapid prototyping and custom internal tools",
            "   • Good for pilot projects",
            "",
            "3. Code-First Frameworks",
            "   • LangChain, LlamaIndex, Semantic Kernel, AutoGen",
            "   • Maximum flexibility and control",
            "   • Requires strong engineering team (3+ ML engineers)",
            "",
            "4. Specialized Tools",
            "   • Document AI, Textract (OCR), Drools (rules engine)",
            "   • Point solutions for specific capabilities"
        ]
    )

    # Slide 5: Top Enterprise Platforms
    add_table_slide(
        "Enterprise Platforms Comparison",
        ["Platform", "Best For", "Compliance", "Est. Cost/Year", "Rating"],
        [
            ["Pega", "Case management", "SOC 2, ISO", "$150K-300K+", "★★★★★"],
            ["Appian", "Rapid dev, govt", "FedRAMP, SOC 2", "$100K-250K", "★★★★★"],
            ["Azure AI Studio", "Microsoft shops", "FedRAMP High", "$75K-200K", "★★★★★"],
            ["Vertex AI", "Google shops", "FedRAMP High", "$50K-150K", "★★★★☆"],
            ["AWS Bedrock", "AWS shops", "FedRAMP High", "$50K-150K", "★★★★☆"],
            ["UiPath AI Center", "Legacy systems", "SOC 2", "$50K-200K", "★★★★★"],
            ["ServiceNow", "IT workflows", "FedRAMP Mod", "$50K-100K", "★★★★☆"]
        ]
    )

    # Slide 6: Platform Deep Dive - Pega
    add_content_slide(
        "Spotlight: Pega Platform",
        [
            "Why Pega for Benefits Adjudication?",
            "",
            "✓ Purpose-built for complex case management",
            "   • Industry-leading case management capabilities",
            "   • Multi-step workflow orchestration",
            "   • Built-in appeal management",
            "",
            "✓ Decision automation with rules engine",
            "   • Handle complex eligibility criteria",
            "   • Decision tables and business rules",
            "",
            "✓ Government/public sector proven",
            "   • Existing customer base in social services",
            "   • Compliance certifications",
            "",
            "⚠ Considerations: Higher cost ($150K-300K/year), implementation complexity"
        ]
    )

    # Slide 7: Platform Deep Dive - Azure
    add_content_slide(
        "Spotlight: Microsoft Azure AI Studio",
        [
            "Why Azure for Benefits Adjudication?",
            "",
            "✓ Best for Microsoft-based government agencies",
            "   • Integration with Office 365, Teams, Power Platform",
            "   • Azure Government Cloud (FedRAMP High)",
            "",
            "✓ Copilot Studio + Power Automate",
            "   • Low-code workflow automation",
            "   • Custom AI agents with enterprise security",
            "",
            "✓ Azure OpenAI Service",
            "   • GPT-4, Claude integration",
            "   • Content safety and responsible AI tools",
            "",
            "✓ Azure Forms Recognizer",
            "   • Document processing for claims and forms",
            "",
            "⚠ Considerations: Requires Azure expertise, consumption-based pricing complexity"
        ]
    )

    # Slide 8: Platform Deep Dive - UiPath
    add_content_slide(
        "Spotlight: UiPath AI Center",
        [
            "Why UiPath for Benefits Adjudication?",
            "",
            "✓ Excellent for legacy system integration",
            "   • RPA bots interact with systems without APIs",
            "   • No need to modernize legacy infrastructure immediately",
            "",
            "✓ Document Understanding AI",
            "   • Extract data from unstructured documents",
            "   • Medical records, claims forms, supporting docs",
            "",
            "✓ Action Center for human-in-the-loop",
            "   • Route edge cases to human adjudicators",
            "   • Audit trail of human decisions",
            "",
            "✓ Process mining and analytics",
            "   • Identify bottlenecks and optimization opportunities",
            "",
            "⚠ Considerations: Platform licensing costs, RPA maintenance overhead"
        ]
    )

    # Slide 9: Low-Code Options
    add_content_slide(
        "Low-Code Platforms for Prototyping",
        [
            "Best for: Rapid pilots and internal tools",
            "",
            "Retool AI - ★★★★☆",
            "• Build custom admin interfaces for case review",
            "• Connect to databases and APIs easily",
            "• $10-50/user/month, SOC 2 compliant",
            "",
            "Flowise / LangFlow - ★★★☆☆",
            "• Visual agent builder (LangChain UI)",
            "• Free/open-source, self-hosted",
            "• Good for prototyping, requires engineering for production",
            "",
            "n8n - ★★★☆☆",
            "• Workflow automation with 400+ integrations",
            "• Self-hosted option for data security",
            "• Good for connecting systems, not complex decision logic"
        ]
    )

    # Slide 10: Code Frameworks
    add_content_slide(
        "Code-First Frameworks (Build Your Own)",
        [
            "Best for: Maximum customization and control",
            "",
            "LangChain / LangGraph - ★★★★☆",
            "• Most mature ecosystem for AI agents",
            "• Stateful workflows, multi-agent coordination",
            "• Requires 3+ ML engineers, $2M+ over 3 years",
            "",
            "LlamaIndex - ★★★★☆",
            "• Specialized in document retrieval (RAG)",
            "• Excellent for policy document search",
            "",
            "Microsoft Semantic Kernel - ★★★★☆",
            "• Best for .NET government systems",
            "• Enterprise patterns, Azure integration",
            "",
            "When to build: Highly custom rules, strong eng team, long-term optimization"
        ]
    )

    # Slide 11: Buy vs Build Decision Matrix
    add_table_slide(
        "Buy vs. Build Decision Matrix",
        ["Factor", "Platform/SaaS", "Custom Build"],
        [
            ["Time to Value", "✓✓✓ Fast (3-6 mo)", "✗ Slow (9-18 mo)"],
            ["Total Cost (3yr)", "$$$ $1M-2M", "$$$$ $2.3M-3.2M"],
            ["Compliance", "✓✓ Built-in", "✓ Your responsibility"],
            ["Customization", "✓ Config limits", "✓✓✓ Unlimited"],
            ["Vendor Lock-in", "✗✗ High risk", "✓ Full control"],
            ["Maintenance", "✓✓ Vendor-managed", "✗✗ Your team"],
            ["Skills Required", "✓ Less technical", "✗✗ ML engineers"],
            ["Explainability", "~ Varies", "✓✓ Full control"]
        ]
    )

    # Slide 12: Cost Analysis
    add_content_slide(
        "Total Cost of Ownership (3 Years)",
        [
            "Enterprise Platform (Pega, Appian):",
            "• Licensing: $450K-900K",
            "• Implementation: $200K-500K (year 1)",
            "• Ongoing support: $200K (years 2-3)",
            "• Infrastructure: $150K",
            "• Total: $1M-2M",
            "",
            "Cloud Platform (Azure, Vertex, Bedrock):",
            "• Platform usage: $150K-450K",
            "• Development: $500K-700K",
            "• Engineering (2 FTE): $900K",
            "• Infrastructure: $225K",
            "• Total: $1.8M-2.3M",
            "",
            "Custom Build (LangChain, open-source):",
            "• Engineering (3-5 FTE): $1.2M-2.1M",
            "• Infrastructure: $300K",
            "• LLM APIs: $225K",
            "• Initial development: $500K",
            "• Total: $2.3M-3.2M"
        ]
    )

    # Slide 13: Risk Assessment
    add_table_slide(
        "Risk Comparison",
        ["Risk", "Platform/SaaS", "Custom Build", "Mitigation"],
        [
            ["Vendor discontinue", "High", "N/A", "Choose established vendors"],
            ["Security breach", "Medium", "Medium", "FedRAMP compliance, audits"],
            ["Cost overruns", "Medium", "High", "Fixed licensing vs. dev"],
            ["Failed implementation", "Low-Med", "High", "Platforms more proven"],
            ["Lock-in", "High", "Low", "Hybrid approach"],
            ["Compliance gaps", "Low*", "Medium", "*If FedRAMP certified"],
            ["Skills shortage", "Low", "High", "Platform = config skills"]
        ]
    )

    # Slide 14: Recommended Approach
    add_content_slide(
        "Recommended: Hybrid Approach",
        [
            "Platform for Foundation + Custom AI for Decisions",
            "",
            "Choose Primary Platform based on current tech stack:",
            "• Microsoft-heavy → Azure AI Studio + Copilot Studio",
            "• Google Cloud → Vertex AI Agent Builder",
            "• Best case management → Pega (higher cost, purpose-built)",
            "• Rapid development → Appian",
            "• Heavy legacy systems → UiPath AI Center",
            "",
            "Supplement with Custom Components:",
            "• Document processing: Google Document AI or Azure Forms Recognizer",
            "• Decision agents: LangChain + Anthropic Claude (complex reasoning)",
            "• Rules engine: Drools or platform's native rules",
            "• Observability: LangSmith or platform monitoring",
            "",
            "Why Hybrid? Balance speed-to-market, compliance, and customization"
        ]
    )

    # Slide 15: Example Architecture
    add_content_slide(
        "Reference Architecture (Hybrid)",
        [
            "Citizen Portal (Appian/Pega UI)",
            "        ↓",
            "Workflow Orchestration (Platform)",
            "        ↓",
            "├─ Document Processing (Google Document AI / Azure)",
            "│",
            "├─ Custom AI Agent (LangChain + Anthropic Claude)",
            "│   ├─ RAG for policy retrieval (Pinecone + LlamaIndex)",
            "│   ├─ Rules Engine (Drools) for eligibility",
            "│   └─ Decision explainer",
            "│",
            "└─ Case Management (Appian/Pega)",
            "        ↓",
            "Human Review (for edge cases)",
            "        ↓",
            "Decision Output + Audit Log"
        ]
    )

    # Slide 16: Implementation Roadmap
    add_content_slide(
        "Phased Implementation Roadmap",
        [
            "Phase 1: Pilot (Months 1-3) - $10K-50K",
            "• Select 1-2 simple benefit types (e.g., unemployment)",
            "• Low-code platform for rapid prototyping",
            "• Test with historical cases, measure accuracy",
            "",
            "Phase 2: Production MVP (Months 4-9) - $100K-300K",
            "• Production platform selection based on pilot",
            "• Build integrations (case mgmt, databases)",
            "• Human-in-the-loop workflows",
            "• Security/compliance review, staff training",
            "",
            "Phase 3: Scale (Months 10-18) - $200K-500K/year",
            "• Expand to additional benefit types",
            "• Add complex rule scenarios",
            "• Continuous monitoring and improvement"
        ]
    )

    # Slide 17: Success Metrics
    add_content_slide(
        "Pilot Success Metrics",
        [
            "Measure These in 3-Month Pilot:",
            "",
            "✓ Accuracy",
            "   • >90% agreement with human adjudicators on pilot cases",
            "",
            "✓ Speed",
            "   • 50%+ reduction in time-to-decision for routine cases",
            "",
            "✓ Explainability",
            "   • 100% of decisions have human-readable explanations",
            "",
            "✓ Cost",
            "   • Pilot delivered for <$50K",
            "",
            "✓ User Acceptance",
            "   • Adjudicators find the tool helpful (surveys, feedback)"
        ]
    )

    # Slide 18: Vendor Evaluation Checklist
    add_content_slide(
        "Vendor Demo Evaluation Criteria",
        [
            "During vendor demos, verify:",
            "",
            "□ Integration with our case management system?",
            "□ Support for complex, multi-step eligibility rules?",
            "□ Decision explainability for audit/appeals?",
            "□ FedRAMP authorized or path to authorization?",
            "□ Handle our document types (medical, forms, etc.)?",
            "□ Human-in-the-loop for edge cases?",
            "□ Total cost: licensing + implementation + ongoing?",
            "□ Realistic timeline to production?",
            "□ Customization vs. configuration limits?",
            "□ Data portability (avoid lock-in)?",
            "□ References from government agencies?"
        ]
    )

    # Slide 19: Next Steps
    add_content_slide(
        "Immediate Next Steps (Weeks 1-4)",
        [
            "Week 1-2: Internal Alignment",
            "• Review this research with stakeholders",
            "• Define 2-3 benefit types for pilot",
            "• Document current systems and integration points",
            "• Establish budget authority",
            "",
            "Week 3-4: Vendor Engagement",
            "• Schedule demos with top 5 platforms:",
            "  - Pega, Appian, Azure AI Studio, Vertex AI, UiPath",
            "• Prepare RFI with requirements",
            "• Share anonymized sample data for demos",
            "",
            "Week 5-8: Proof of Concept",
            "• Select top 2 platforms",
            "• Run PoC with real (anonymized) cases",
            "• Measure accuracy, speed, user feedback"
        ]
    )

    # Slide 20: Conclusion
    add_content_slide(
        "Key Takeaways",
        [
            "1. Recommendation: Hybrid Platform Approach",
            "   • Enterprise platform for workflow + case management",
            "   • Custom AI components for complex decision logic",
            "",
            "2. Top Platform Choices:",
            "   • Pega (best case mgmt), Appian (rapid dev),",
            "   • Azure (Microsoft shops), UiPath (legacy systems)",
            "",
            "3. Investment: $1M-2M over 3 years (platform-based)",
            "",
            "4. Timeline: 3-6 months pilot, 9-12 months to production",
            "",
            "5. Critical Success Factors:",
            "   • Start small (1-2 benefit types)",
            "   • Ensure explainability from day one",
            "   • Plan for human oversight (not 100% automation)",
            "   • Choose FedRAMP-authorized platforms"
        ]
    )

    # Save presentation
    output_path = "/home/user/AI-powered-PM/AI decision agents/AI_Agent_Platform_Research_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
